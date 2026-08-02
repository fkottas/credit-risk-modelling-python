"""Characteristic-analysis exports suitable for review packs and teaching."""

from __future__ import annotations

import html
from pathlib import Path

import pandas as pd

from .model import LogisticScorecard


def characteristic_summary(scorecard: LogisticScorecard) -> pd.DataFrame:
    points = scorecard.points_table()
    return (
        points.groupby("feature", as_index=False)
        .agg(
            bins=("bin", "count"),
            observations=("count", "sum"),
            information_value=("iv_component", "sum"),
            min_bad_rate=("bad_rate", "min"),
            max_bad_rate=("bad_rate", "max"),
            points_range=("points", lambda values: float(values.max() - values.min())),
        )
        .sort_values("information_value", ascending=False)
    )


def export_characteristic_report(
    scorecard: LogisticScorecard, output_dir: str | Path
) -> dict[str, Path]:
    """Write CSV tables and a dependency-light HTML review pack."""

    destination = Path(output_dir)
    destination.mkdir(parents=True, exist_ok=True)
    points = scorecard.points_table()
    summary = characteristic_summary(scorecard)
    points_path = destination / "scorecard_points.csv"
    summary_path = destination / "characteristic_summary.csv"
    html_path = destination / "characteristic_analysis.html"
    points.to_csv(points_path, index=False)
    summary.to_csv(summary_path, index=False)
    sections = []
    for feature, table in points.groupby("feature", sort=False):
        sections.append(
            f"<h2>{html.escape(str(feature))}</h2>{table.to_html(index=False, float_format=lambda x: f'{x:.5f}')}"
        )
    document = f"""<!doctype html>
<html lang="en"><head><meta charset="utf-8"><title>Characteristic analysis</title>
<style>body{{font:14px Arial;max-width:1200px;margin:2rem auto;color:#18324b}}table{{border-collapse:collapse;width:100%;margin-bottom:2rem}}th,td{{border:1px solid #c8d4df;padding:.4rem;text-align:right}}th{{background:#e8eef5}}td:first-child,td:nth-child(2){{text-align:left}}h1,h2{{color:#1f4d78}}</style>
</head><body><h1>Scorecard characteristic analysis</h1><p>Generated from the fitted training sample. Validate independently before use.</p>
<h2>Summary</h2>{summary.to_html(index=False, float_format=lambda x: f"{x:.5f}")}{"".join(sections)}</body></html>"""
    html_path.write_text(document, encoding="utf-8")
    return {"points": points_path, "summary": summary_path, "html": html_path}


def export_characteristic_presentation(
    scorecard: LogisticScorecard,
    output_path: str | Path,
    *,
    title: str = "Scorecard characteristic analysis",
) -> Path:
    """Create a review-ready PowerPoint with one slide per characteristic.

    ``python-pptx`` is an optional book dependency.  The presentation is an
    evidence pack, not automatic approval: the modeller must add sample dates,
    population definitions, policy decisions, validation conclusions, and
    sign-offs before institutional use.
    """

    try:
        from pptx import Presentation
        from pptx.chart.data import ChartData
        from pptx.dml.color import RGBColor
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as error:  # pragma: no cover - exercised when optional extra is absent
        raise ImportError("Install creditriskbook[book] to export PowerPoint reports") from error

    destination = Path(output_path)
    destination.parent.mkdir(parents=True, exist_ok=True)
    presentation = Presentation()
    presentation.slide_width, presentation.slide_height = Inches(13.333), Inches(7.5)
    navy, blue, pale = RGBColor(24, 50, 75), RGBColor(31, 77, 120), RGBColor(232, 238, 245)

    cover = presentation.slides.add_slide(presentation.slide_layouts[0])
    cover.shapes.title.text = title
    cover.placeholders[1].text = (
        "Binning, bad rates, WOE, information value and score points\n"
        "Generated from the fitted development sample — independent validation required"
    )
    for shape in cover.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Aptos"
                    run.font.color.rgb = navy

    summary = characteristic_summary(scorecard)
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Characteristic summary"
    rows, columns = len(summary) + 1, len(summary.columns)
    table = slide.shapes.add_table(
        rows, columns, Inches(0.45), Inches(1.25), Inches(12.4), Inches(5.6)
    ).table
    for column, value in enumerate(summary.columns):
        table.cell(0, column).text = str(value).replace("_", " ").title()
    for row, record in enumerate(summary.itertuples(index=False), start=1):
        for column, value in enumerate(record):
            table.cell(row, column).text = (
                f"{value:.4f}" if isinstance(value, float) else str(value)
            )
    for row in range(rows):
        for column in range(columns):
            cell = table.cell(row, column)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                blue if row == 0 else (pale if row % 2 else RGBColor(255, 255, 255))
            )
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name, run.font.size = "Aptos", Pt(8.5)
                    run.font.bold = row == 0
                    run.font.color.rgb = RGBColor(255, 255, 255) if row == 0 else navy

    points = scorecard.points_table()
    for feature, characteristic in points.groupby("feature", sort=False):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = f"Characteristic: {feature}"
        chart_data = ChartData()
        labels = [str(value)[:28] for value in characteristic["bin"]]
        chart_data.categories = labels
        chart_data.add_series("Bad rate", characteristic["bad_rate"].astype(float).tolist())
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.55),
            Inches(1.35),
            Inches(6.0),
            Inches(4.7),
            chart_data,
        ).chart
        chart.has_legend = False
        chart.value_axis.maximum_scale = max(
            0.10, min(1.0, float(characteristic["bad_rate"].max()) * 1.25)
        )
        chart.value_axis.has_major_gridlines = True
        visible = characteristic[["bin", "count", "bad_rate", "woe", "points"]]
        shape = slide.shapes.add_table(
            len(visible) + 1,
            len(visible.columns),
            Inches(6.75),
            Inches(1.35),
            Inches(6.0),
            Inches(4.7),
        )
        detail_table = shape.table
        for column, value in enumerate(visible.columns):
            detail_table.cell(0, column).text = str(value).replace("_", " ").title()
        for row, record in enumerate(visible.itertuples(index=False), start=1):
            for column, value in enumerate(record):
                detail_table.cell(row, column).text = (
                    f"{value:.4f}" if isinstance(value, float) else str(value)[:34]
                )
        note = slide.shapes.add_textbox(Inches(0.6), Inches(6.25), Inches(12.0), Inches(0.55))
        note.text_frame.text = (
            f"IV={characteristic['iv_component'].sum():.4f}. Challenge cut points, sparse bins, "
            "time stability, leakage, business meaning and adverse-action suitability."
        )
        for paragraph in note.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name, run.font.size, run.font.color.rgb = "Aptos", Pt(9), navy
    presentation.save(destination)
    return destination
